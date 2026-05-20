"""
core/extractor.py -- Feature extraction from files.

FileFeatures is the shared contract between the extractor and all classifiers.
Text is capped at 4000 characters to keep TF-IDF fast.

Supported types:  .pdf  .docx  .txt  .zip
Optional (PPTX):  .pptx  (imported defensively; falls back to filename-only)

For any unsupported type or any extraction failure, text is an empty string.
The rest of the pipeline handles empty text gracefully (filename-only mode).
"""

import os
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

TEXT_CAP = 4000  # max characters extracted from document content


@dataclass
class FileFeatures:
    """All signals extracted from a single file."""
    path: str
    filename: str           # e.g. "cs180_lab3.pdf"
    stem: str               # e.g. "cs180_lab3"
    ext: str                # lowercase, with dot, e.g. ".pdf"
    size_bytes: int
    text: str               # extracted content, up to TEXT_CAP chars; "" on failure
    zip_members: list[str] = field(default_factory=list)   # for .zip files

    @property
    def all_text(self) -> str:
        """Combine filename stem + extracted text + zip member names for classification."""
        parts = [self.stem.replace("_", " ").replace("-", " ")]
        if self.text:
            parts.append(self.text)
        if self.zip_members:
            parts.append(" ".join(m.replace("_", " ") for m in self.zip_members))
        return " ".join(parts)


def extract_features(path: str) -> FileFeatures:
    """
    Extract features from a file at `path`.
    Never raises — returns FileFeatures with empty text/members on any failure.
    """
    p = Path(path)
    ext = p.suffix.lower()
    size = 0
    try:
        size = p.stat().st_size
    except OSError:
        pass

    text = ""
    zip_members: list[str] = []

    try:
        if ext == ".pdf":
            text = _extract_pdf(path)
        elif ext == ".docx":
            text = _extract_docx(path)
        elif ext == ".pptx":
            text = _extract_pptx(path)
        elif ext == ".txt":
            text = _extract_txt(path)
        elif ext == ".zip":
            zip_members = _inspect_zip(path)
    except Exception:
        # Extraction failure is non-fatal — fall back to filename-only.
        text = ""
        zip_members = []

    return FileFeatures(
        path=str(p),
        filename=p.name,
        stem=p.stem,
        ext=ext,
        size_bytes=size,
        text=text[:TEXT_CAP],
        zip_members=zip_members,
    )


def extract_folder_features(path: str) -> FileFeatures:
    """
    Extract features from a directory.
    Uses folder name + contained filenames (as zip_members) + limited text
    from the first few supported files found inside.
    Never raises.
    """
    p = Path(path)
    total_size = 0
    member_names: list[str] = []
    text_parts: list[str] = []

    try:
        for f in sorted(p.rglob("*")):
            if not f.is_file():
                continue
            try:
                total_size += f.stat().st_size
            except OSError:
                pass
            member_names.append(f.name)
            if len(text_parts) < 3 and f.suffix.lower() in {".pdf", ".docx", ".txt", ".pptx"}:
                try:
                    snippet = extract_features(str(f)).text
                    if snippet:
                        text_parts.append(snippet[:500])
                except Exception:
                    pass
    except Exception:
        pass

    return FileFeatures(
        path=str(p),
        filename=p.name,
        stem=p.name,
        ext="",
        size_bytes=total_size,
        text=" ".join(text_parts)[:TEXT_CAP],
        zip_members=member_names[:50],
    )


# ---------------------------------------------------------------------------
# Per-format helpers
# ---------------------------------------------------------------------------

def _extract_pdf(path: str) -> str:
    """Try pdfplumber first, fall back to PyMuPDF (fitz)."""
    text = _pdf_via_pdfplumber(path)
    if not text.strip():
        text = _pdf_via_fitz(path)
    return text[:TEXT_CAP]


def _pdf_via_pdfplumber(path: str) -> str:
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(path) as pdf:
            parts = []
            for page in pdf.pages[:10]:  # first 10 pages is enough for signals
                page_text = page.extract_text() or ""
                parts.append(page_text)
                if sum(len(p) for p in parts) >= TEXT_CAP:
                    break
            return " ".join(parts)
    except Exception:
        return ""


def _pdf_via_fitz(path: str) -> str:
    try:
        import fitz  # PyMuPDF  # type: ignore
        doc = fitz.open(path)
        parts = []
        for page in doc[:10]:
            parts.append(page.get_text())
            if sum(len(p) for p in parts) >= TEXT_CAP:
                break
        doc.close()
        return " ".join(parts)
    except Exception:
        return ""


def _extract_docx(path: str) -> str:
    try:
        from docx import Document  # type: ignore
        doc = Document(path)
        text = " ".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:TEXT_CAP]
    except Exception:
        return ""


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
                if sum(len(p) for p in parts) >= TEXT_CAP:
                    break
        return " ".join(parts)[:TEXT_CAP]
    except Exception:
        return ""


def _extract_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(TEXT_CAP)
    except Exception:
        return ""


def _inspect_zip(path: str) -> list[str]:
    """Return the list of member filenames inside a ZIP archive."""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            # Return only the filename component, not full paths inside the zip.
            names = []
            for info in zf.infolist():
                if not info.is_dir():
                    names.append(Path(info.filename).name)
            return names[:50]  # cap to 50 members
    except Exception:
        return []
